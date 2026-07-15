"""Document ingestion pipeline — extract, chunk, embed."""

from __future__ import annotations

import csv
import io
import logging
import math
from pathlib import Path

from django.core.files.storage import default_storage

from apps.knowledge.infrastructure.models import KnowledgeChunk, KnowledgeDocument
from infrastructure.vector.embeddings import EmbeddingService, get_embedding_service

logger = logging.getLogger("apps.knowledge")

CHUNK_SIZE = 800
CHUNK_OVERLAP = 100


class DocumentIngestionService:
    """Detect type → extract text → chunk → embed → store."""

    def __init__(self, embedding_service: EmbeddingService | None = None) -> None:
        self._embeddings = embedding_service or get_embedding_service()

    def process(self, document_id: str) -> KnowledgeDocument:
        doc = KnowledgeDocument.objects.get(pk=document_id)
        doc.status = KnowledgeDocument.Status.PROCESSING
        doc.error_message = ""
        doc.save(update_fields=["status", "error_message", "updated_at"])

        try:
            file_type = self.detect_type(doc)
            doc.file_type = file_type
            text, meta = self.extract_text(doc, file_type)
            chunks = self.chunk_text(text)
            embeddings = self._embeddings.embed_many(chunks) if chunks else []

            KnowledgeChunk.objects.filter(document=doc).delete()
            for idx, (content, embedding) in enumerate(zip(chunks, embeddings, strict=True)):
                KnowledgeChunk.objects.create(
                    document=doc,
                    index=idx,
                    content=content,
                    embedding=embedding,
                    metadata={"page": meta.get("pages", 1)},
                )

            doc.page_count = int(meta.get("pages", 1) or 1)
            doc.chunk_count = len(chunks)
            doc.metadata = {**(doc.metadata or {}), **meta}
            doc.status = KnowledgeDocument.Status.READY
            doc.save(
                update_fields=[
                    "file_type",
                    "page_count",
                    "chunk_count",
                    "metadata",
                    "status",
                    "updated_at",
                ]
            )
        except Exception as exc:
            logger.exception("document_ingestion_failed", extra={"document_id": document_id})
            doc.status = KnowledgeDocument.Status.FAILED
            doc.error_message = str(exc)[:2000]
            doc.save(update_fields=["status", "error_message", "updated_at"])

        return doc

    def detect_type(self, doc: KnowledgeDocument) -> str:
        name = ""
        if doc.file:
            name = Path(doc.file.name).name.lower()
        elif doc.title:
            name = doc.title.lower()
        ext = Path(name).suffix.lstrip(".")
        mapping = {
            "pdf": KnowledgeDocument.FileType.PDF,
            "docx": KnowledgeDocument.FileType.DOCX,
            "doc": KnowledgeDocument.FileType.DOCX,
            "pptx": KnowledgeDocument.FileType.PPTX,
            "ppt": KnowledgeDocument.FileType.PPTX,
            "xlsx": KnowledgeDocument.FileType.XLSX,
            "xls": KnowledgeDocument.FileType.XLSX,
            "csv": KnowledgeDocument.FileType.CSV,
            "png": KnowledgeDocument.FileType.IMAGE,
            "jpg": KnowledgeDocument.FileType.IMAGE,
            "jpeg": KnowledgeDocument.FileType.IMAGE,
            "gif": KnowledgeDocument.FileType.IMAGE,
            "webp": KnowledgeDocument.FileType.IMAGE,
            "tiff": KnowledgeDocument.FileType.IMAGE,
            "bmp": KnowledgeDocument.FileType.IMAGE,
        }
        return mapping.get(ext, KnowledgeDocument.FileType.OTHER)

    def extract_text(self, doc: KnowledgeDocument, file_type: str) -> tuple[str, dict]:
        if not doc.file:
            return f"Document: {doc.title}", {"pages": 1, "source": "title_only"}

        with default_storage.open(doc.file.name, "rb") as fh:
            raw = fh.read()

        extractors = {
            KnowledgeDocument.FileType.PDF: self._extract_pdf,
            KnowledgeDocument.FileType.DOCX: self._extract_docx,
            KnowledgeDocument.FileType.PPTX: self._extract_pptx,
            KnowledgeDocument.FileType.XLSX: self._extract_xlsx,
            KnowledgeDocument.FileType.CSV: self._extract_csv,
            KnowledgeDocument.FileType.IMAGE: self._extract_image,
        }
        extractor = extractors.get(file_type, self._extract_fallback)
        return extractor(raw, doc)

    def chunk_text(
        self, text: str, size: int = CHUNK_SIZE, overlap: int = CHUNK_OVERLAP
    ) -> list[str]:
        text = (text or "").strip()
        if not text:
            return []
        if len(text) <= size:
            return [text]
        chunks: list[str] = []
        start = 0
        step = max(1, size - overlap)
        while start < len(text):
            end = min(len(text), start + size)
            chunk = text[start:end].strip()
            if chunk:
                chunks.append(chunk)
            if end >= len(text):
                break
            start += step
        return chunks

    def _extract_pdf(self, raw: bytes, doc: KnowledgeDocument) -> tuple[str, dict]:
        try:
            from pypdf import PdfReader

            reader = PdfReader(io.BytesIO(raw))
            pages = [page.extract_text() or "" for page in reader.pages]
            text = "\n\n".join(pages).strip()
            return text or f"[PDF with no extractable text] {doc.title}", {
                "pages": len(reader.pages),
                "parser": "pypdf",
            }
        except Exception as exc:
            logger.warning("pdf_extract_failed: %s", exc)
            return f"[PDF parse failed] {doc.title}", {"pages": 1, "error": str(exc)}

    def _extract_docx(self, raw: bytes, doc: KnowledgeDocument) -> tuple[str, dict]:
        try:
            from docx import Document as DocxDocument

            document = DocxDocument(io.BytesIO(raw))
            paras = [p.text for p in document.paragraphs if p.text.strip()]
            return "\n".join(paras) or f"[Empty DOCX] {doc.title}", {
                "pages": max(1, math.ceil(len(paras) / 40)),
                "parser": "python-docx",
            }
        except Exception as exc:
            return f"[DOCX parse failed] {doc.title}", {"pages": 1, "error": str(exc)}

    def _extract_pptx(self, raw: bytes, doc: KnowledgeDocument) -> tuple[str, dict]:
        try:
            from pptx import Presentation

            prs = Presentation(io.BytesIO(raw))
            texts: list[str] = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        texts.append(shape.text)
            return "\n".join(texts) or f"[Empty PPTX] {doc.title}", {
                "pages": len(prs.slides),
                "parser": "python-pptx",
            }
        except Exception as exc:
            return f"[PPTX parse failed] {doc.title}", {"pages": 1, "error": str(exc)}

    def _extract_xlsx(self, raw: bytes, doc: KnowledgeDocument) -> tuple[str, dict]:
        try:
            from openpyxl import load_workbook

            wb = load_workbook(io.BytesIO(raw), read_only=True, data_only=True)
            rows: list[str] = []
            for sheet in wb.worksheets:
                rows.append(f"# Sheet: {sheet.title}")
                for row in sheet.iter_rows(values_only=True):
                    cells = [str(c) if c is not None else "" for c in row]
                    if any(cells):
                        rows.append("\t".join(cells))
            return "\n".join(rows) or f"[Empty XLSX] {doc.title}", {
                "pages": len(wb.worksheets),
                "parser": "openpyxl",
            }
        except Exception as exc:
            return f"[XLSX parse failed] {doc.title}", {"pages": 1, "error": str(exc)}

    def _extract_csv(self, raw: bytes, doc: KnowledgeDocument) -> tuple[str, dict]:
        try:
            text = raw.decode("utf-8", errors="replace")
            reader = csv.reader(io.StringIO(text))
            lines = [", ".join(row) for row in reader]
            return "\n".join(lines) or f"[Empty CSV] {doc.title}", {
                "pages": 1,
                "parser": "csv",
            }
        except Exception as exc:
            return f"[CSV parse failed] {doc.title}", {"pages": 1, "error": str(exc)}

    def _extract_image(self, raw: bytes, doc: KnowledgeDocument) -> tuple[str, dict]:
        try:
            import pytesseract
            from PIL import Image

            image = Image.open(io.BytesIO(raw))
            text = pytesseract.image_to_string(image)
            if text.strip():
                return text.strip(), {"pages": 1, "parser": "pytesseract", "ocr": True}
            return (
                f"[OCR empty] {doc.title}",
                {"pages": 1, "parser": "pytesseract", "ocr": True, "ocr_empty": True},
            )
        except Exception as exc:
            logger.info("ocr_unavailable: %s", exc)
            filename = Path(doc.file.name).name if doc.file else doc.title
            return (
                f"[OCR unavailable] Filename: {filename}. Title: {doc.title}.",
                {
                    "pages": 1,
                    "parser": "stub",
                    "ocr": False,
                    "ocr_error": str(exc),
                    "filename": filename,
                },
            )

    def _extract_fallback(self, raw: bytes, doc: KnowledgeDocument) -> tuple[str, dict]:
        try:
            text = raw.decode("utf-8", errors="replace").strip()
            if text:
                return text, {"pages": 1, "parser": "text"}
        except Exception:
            pass
        return f"[Unsupported file] {doc.title}", {"pages": 1, "parser": "stub"}
